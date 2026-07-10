"""第 6 版周报 PPT：
- Pages 1-2 保持不动（头条 + 创新点）
- Pages 3-8 图片大幅放大，每页只放 1-2 张图
- val_batch 已裁剪成 2-检测/张，放到 docs/crops/
- 失败案例单图占满页
"""
from pathlib import Path
from tempfile import NamedTemporaryFile
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image, ImageDraw

ROOT = Path(__file__).parent.parent
RUN_DIR = ROOT / 'runs' / 'coil_loss_ablation' / '02_coverage_v4_N_model'
FN_DIR = RUN_DIR / 'fn_samples'
CROP_DIR = ROOT / 'docs' / 'crops'
OUT = ROOT / 'docs' / 'weekly_summary.pptx'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

NAVY = RGBColor(0x1E, 0x40, 0xAF)
NAVY_DARK = RGBColor(0x1E, 0x3A, 0x8A)
GREEN = RGBColor(0x16, 0x6E, 0x3A)
ORANGE = RGBColor(0xEA, 0x58, 0x0C)
GRAY = RGBColor(0x6B, 0x72, 0x80)
GRAY_LIGHT = RGBColor(0xE5, 0xE7, 0xEB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xDC, 0x26, 0x26)
BG_LIGHT = RGBColor(0xF9, 0xFA, 0xFB)


def add_text(slide, left, top, width, height, text, *,
             size=18, bold=False, color=NAVY_DARK, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = 'Microsoft YaHei'
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_rect(slide, left, top, width, height, fill=NAVY, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_image_fit(slide, img_path, left, top, max_w, max_h, max_px=None):
    if max_px is not None:
        with Image.open(img_path) as im:
            iw, ih = im.size
            big = max(iw, ih)
            if big > max_px:
                scale = max_px / big
                im = im.resize((int(iw * scale), int(ih * scale)),
                               Image.LANCZOS)
            tmp = NamedTemporaryFile(suffix='.jpg', delete=False)
            im.convert('RGB').save(tmp.name, 'JPEG', quality=92)
            tmp.close()
            use_path = tmp.name
    else:
        use_path = str(img_path)
    with Image.open(use_path) as im:
        iw, ih = im.size
    ratio = min(max_w / iw, max_h / ih)
    w, h = int(iw * ratio), int(ih * ratio)
    x = int(left + (max_w - w) / 2)
    y = int(top + (max_h - h) / 2)
    slide.shapes.add_picture(use_path, x, y, w, h)
    return (x, y, w, h)


def add_header(slide, title, subtitle):
    add_text(slide, Inches(0.5), Inches(0.3), Inches(11), Inches(0.55),
             title, size=24, bold=True, color=NAVY_DARK)
    add_text(slide, Inches(0.5), Inches(0.85), Inches(11), Inches(0.35),
             subtitle, size=11, color=GRAY)
    add_rect(slide, Inches(0.5), Inches(1.25), Inches(12.3), Inches(0.04),
             fill=NAVY)


def add_pagenum(slide, n, total=8):
    add_text(slide, Inches(12.4), Inches(7.15), Inches(0.7), Inches(0.25),
             f'{n} / {total}', size=10, color=GRAY, align=PP_ALIGN.RIGHT)


def make_coverage_diagram():
    W, H = 1200, 600
    im = Image.new('RGB', (W, H), WHITE)
    d = ImageDraw.Draw(im)
    box_top = 80
    d.text((50, 30), 'IoU 视角', fill=NAVY_DARK)
    d.rectangle([(80, box_top), (480, box_top + 360)], outline=GREEN, width=4)
    d.ellipse([275, box_top + 175, 285, box_top + 185], fill=GREEN)
    d.rectangle([(580, box_top + 280), (780, box_top + 380)], outline=RED, width=4)
    d.ellipse([675, box_top + 325, 685, box_top + 335], fill=RED)
    d.text((180, box_top + 380), 'GT（宽松标注）', fill=GREEN)
    d.text((600, box_top + 395), 'Pred', fill=RED)
    d.text((180, box_top + 460), 'IoU ≈ 0.05', fill=GRAY)
    d.text((180, box_top + 490), '→ 训练无效', fill=RED)
    d.text((650, 30), 'Coverage 视角', fill=NAVY_DARK)
    d.rectangle([(680, box_top), (1080, box_top + 360)], outline=GREEN, width=4)
    d.rectangle([(830, box_top + 130), (980, box_top + 230)], outline=RED, width=4)
    d.ellipse([900, box_top + 175, 910, box_top + 185], fill=GREEN)
    d.text((770, box_top + 250), 'GT', fill=GREEN)
    d.text((835, box_top + 245), 'Pred', fill=RED)
    d.text((830, box_top + 460), 'Coverage ≈ 1.0', fill=GRAY)
    d.text((830, box_top + 490), '→ 训练有效', fill=GREEN)
    tmp = NamedTemporaryFile(suffix='.jpg', delete=False)
    im.save(tmp.name, 'JPEG', quality=95)
    tmp.close()
    return tmp.name


def make_shrink_diagram():
    W, H = 1200, 500
    im = Image.new('RGB', (W, H), WHITE)
    d = ImageDraw.Draw(im)
    cx, cy = 200, 240
    base_w, base_h = 200, 200
    scales = [(1.2, ORANGE, '× 1.2'), (1.0, NAVY, '× 1.0'),
              (0.9, GREEN, '× 0.9'), (0.8, RED, '× 0.8')]
    for i, (s, color, label) in enumerate(scales):
        center_x = 200 + i * 280
        w = int(base_w * s); h = int(base_h * s)
        x1 = center_x - w // 2; y1 = cy - h // 2
        x2 = center_x + w // 2; y2 = cy + h // 2
        d.rectangle([(x1, y1), (x2, y2)], outline=color, width=4)
        d.ellipse([center_x - 6, cy - 6, center_x + 6, cy + 6], fill=color)
        d.text((center_x - 30, y2 + 10), label, fill=color)
    d.text((50, 20), 'bbox_random_shrink：保持中心，边长随机缩放',
           fill=NAVY_DARK)
    d.text((50, 50), '同一目标在训练中以 4 种尺寸出现，教会模型容忍形状差异',
           fill=GRAY)
    for i in range(4):
        d.line([(200 + i * 280, 350), (200 + i * 280, 380)],
               fill=GRAY, width=2)
    d.text((150, 395), '中心点不变', fill=GRAY)
    tmp = NamedTemporaryFile(suffix='.jpg', delete=False)
    im.save(tmp.name, 'JPEG', quality=95)
    tmp.close()
    return tmp.name


def make_nwd_diagram():
    W, H = 1200, 600
    im = Image.new('RGB', (W, H), WHITE)
    d = ImageDraw.Draw(im)
    d.text((50, 30), 'Step 1：Bbox → 2D Gaussian', fill=NAVY_DARK)
    gt_box = [(150, 150), (450, 350)]
    d.rectangle(gt_box, outline=GREEN, width=4)
    cx, cy = 300, 250
    for r in [80, 60, 40, 20]:
        d.ellipse([cx - r, cy - r, cx + r, cy + r], outline=GREEN, width=2)
    d.text((150, 360), 'GT bbox → N(μ, Σ)', fill=GREEN)
    d.text((150, 395), 'Σ = diag((w/2)², (h/2)²)', fill=GRAY)
    d.text((430, 30), 'Step 2：Pred 同理建模', fill=NAVY_DARK)
    pred_box = [(650, 200), (820, 340)]
    d.rectangle(pred_box, outline=RED, width=4)
    pcx, pcy = 735, 270
    for r in [70, 50, 30, 15]:
        d.ellipse([pcx - r, pcy - r, pcx + r, pcy + r], outline=RED, width=2)
    d.text((650, 360), 'Pred bbox → N(μ, Σ)', fill=RED)
    d.text((830, 30), 'Step 3：Wasserstein 距离', fill=NAVY_DARK)
    d.text((830, 100), 'W₂²(P, Q) = ||μ₁-μ₂||²', fill=GRAY)
    d.text((830, 130), '           + ||Σ₁^(1/2) - Σ₂^(1/2)||²_F', fill=GRAY)
    d.text((830, 200), 'NWD = exp(-W₂ / C)', fill=NAVY_DARK)
    d.text((830, 250), 'C = 12.0（数据集平均尺度）', fill=GRAY)
    d.text((830, 320), 'Loss = 1 - NWD', fill=RED)
    d.text((830, 380), '位置/形状都对齐 → NWD→1', fill=GREEN)
    tmp = NamedTemporaryFile(suffix='.jpg', delete=False)
    im.save(tmp.name, 'JPEG', quality=95)
    tmp.close()
    return tmp.name


# ============================================================
# SLIDE 1: 标题 + 头条数字（保持不变）
# ============================================================
s1 = prs.slide_layouts[6]
s1 = prs.slides.add_slide(s1)

add_text(s1, Inches(0.5), Inches(0.6), Inches(12.3), Inches(0.7),
         '钢卷头尾小目标检测', size=36, bold=True, color=NAVY_DARK)
add_text(s1, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.4),
         '基于 Hyper-YOLO 的工业场景适配', size=14, color=GRAY)
add_rect(s1, Inches(0.5), Inches(1.85), Inches(12.3), Inches(0.04), fill=NAVY)

add_text(s1, Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.5),
         '学术 mAP@0.5', size=16, color=GRAY, align=PP_ALIGN.CENTER)
add_text(s1, Inches(0.5), Inches(3.0), Inches(12.3), Inches(1.5),
         '0.877', size=120, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(s1, Inches(0.5), Inches(4.6), Inches(12.3), Inches(0.5),
         '相对 baseline +105%', size=18, bold=True, color=GREEN,
         align=PP_ALIGN.CENTER)

add_text(s1, Inches(0.5), Inches(5.4), Inches(12.3), Inches(0.4),
         '部署口径（业务 top1 评估）', size=11, color=GRAY,
         align=PP_ALIGN.CENTER)

card_w = Inches(2.6)
card_gap = Inches(0.4)
total_w = card_w * 3 + card_gap * 2
card_x_start = (SW - total_w) // 2

for i, (name, val) in enumerate([('Recall', '0.868'),
                                  ('Precision', '0.943'),
                                  ('F1', '0.904')]):
    cx = card_x_start + i * (card_w + card_gap)
    add_rect(s1, cx, Inches(5.9), card_w, Inches(1.0), fill=BG_LIGHT)
    add_text(s1, cx, Inches(5.95), card_w, Inches(0.3),
             name, size=10, color=GRAY, align=PP_ALIGN.CENTER)
    add_text(s1, cx, Inches(6.25), card_w, Inches(0.6),
             val, size=28, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

add_pagenum(s1, 1, total=8)


# ============================================================
# SLIDE 2: 创新点详解（保持不变）
# ============================================================
s2 = prs.slide_layouts[6]
s2 = prs.slides.add_slide(s2)
add_header(s2, '核心创新点详解',
           'Coverage Loss + bbox_random_shrink + NWD 三种小目标优化')

left_x = Inches(0.5)
top_y = Inches(1.45)
half_w = Inches(6.2)
half_h = Inches(2.4)

add_text(s2, left_x, top_y, half_w, Inches(0.35),
         '① Coverage Loss：宽松标注下的中心监督',
         size=12, bold=True, color=GREEN)

cov_diag = make_coverage_diagram()
add_rect(s2, left_x, top_y + Inches(0.4), half_w, half_h - Inches(0.4),
         fill=BG_LIGHT)
add_image_fit(s2, cov_diag,
              left_x + Inches(0.05), top_y + Inches(0.42),
              half_w - Inches(0.1), half_h - Inches(0.45),
              max_px=900)

right_x = Inches(6.85)
add_text(s2, right_x, top_y, half_w, Inches(0.35),
         '② bbox_random_shrink：边长随机缩放',
         size=12, bold=True, color=ORANGE)

shrink_diag = make_shrink_diagram()
add_rect(s2, right_x, top_y + Inches(0.4), half_w, half_h - Inches(0.4),
         fill=BG_LIGHT)
add_image_fit(s2, shrink_diag,
              right_x + Inches(0.05), top_y + Inches(0.42),
              half_w - Inches(0.1), half_h - Inches(0.45),
              max_px=900)

bot_y = Inches(4.0)
add_text(s2, Inches(0.5), bot_y, Inches(12.3), Inches(0.35),
         '③ NWD（Normalized Wasserstein Distance）：小目标相似度度量',
         size=12, bold=True, color=NAVY_DARK)

nwd_diag = make_nwd_diagram()
add_rect(s2, Inches(0.5), bot_y + Inches(0.4), Inches(12.3), Inches(2.4),
         fill=BG_LIGHT)
add_image_fit(s2, nwd_diag,
              Inches(0.55), bot_y + Inches(0.42),
              Inches(12.2), Inches(2.35),
              max_px=1100)

add_rect(s2, Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.5),
         fill=NAVY_DARK)
add_text(s2, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4),
         '三种方案在本数据集完全等价（Recall 0.868），任一可作为小目标优化的选择',
         size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)

add_pagenum(s2, 2, total=8)


# ============================================================
# SLIDE 3-5: 成功检测案例（每页 1 张裁剪后的大图）
# ============================================================
success_cases = [
    ('success_1.jpg', '测试图组 1：钢卷缠绕区域的 tip 检测'),
    ('success_3.jpg', '测试图组 2：钢卷表面的 tip 检测'),
    ('success_5.jpg', '测试图组 3：高反光场景下的 tip 检测'),
]

for idx, (img_name, caption) in enumerate(success_cases):
    s = prs.slide_layouts[6]
    s = prs.slides.add_slide(s)
    add_header(s, f'成功检测案例 ({idx+1}/3)',
               'v4 best.pt · imgsz=1024 · conf=0.05 · 红框=模型预测')

    # 单张图占满页（除页眉/页脚）
    img_y = Inches(1.45)
    img_h = Inches(5.4)
    img_w = Inches(12.3)
    img_x = Inches(0.5)

    add_rect(s, img_x, img_y, img_w, img_h, fill=BG_LIGHT)
    add_image_fit(s, CROP_DIR / img_name,
                  img_x + Inches(0.1), img_y + Inches(0.1),
                  img_w - Inches(0.2), img_h - Inches(0.5),
                  max_px=2400)
    add_rect(s, img_x, img_y + img_h - Inches(0.4),
             img_w, Inches(0.35), fill=GREEN)
    add_text(s, img_x, img_y + img_h - Inches(0.4),
             img_w, Inches(0.35),
             f'✓  {caption}', size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_rect(s, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.45),
             fill=NAVY_DARK)
    add_text(s, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.35),
             '测试集 Recall 0.868 / Precision 0.943 / F1 0.904',
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)

    add_pagenum(s, idx + 3, total=8)


# ============================================================
# SLIDE 6: 失败案例 2 张大图（高反光 + 钢丝遮挡）
# ============================================================
s6 = prs.slide_layouts[6]
s6 = prs.slides.add_slide(s6)
add_header(s6, '剩余难例 (1/3)',
           '绿框=真实位置（GT），红框=模型预测（缺失或偏移）')

# 2 张大图 side by side
img_y = Inches(1.45)
img_w = Inches(6.3)
img_h = Inches(5.4)
img_gap = Inches(0.3)
img_x_start = (SW - (img_w * 2 + img_gap)) // 2

fn_pairs = [
    ('262.png', '高反光区', '#1'),
    ('610.png', '钢丝遮挡', '#2'),
]

for i, (fn, label, num) in enumerate(fn_pairs):
    ix = img_x_start + i * (img_w + img_gap)
    add_rect(s6, ix, img_y, img_w, img_h, fill=BG_LIGHT)
    add_image_fit(s6, FN_DIR / fn,
                  ix + Inches(0.1), img_y + Inches(0.1),
                  img_w - Inches(0.2), img_h - Inches(0.5),
                  max_px=2400)
    add_rect(s6, ix, img_y + img_h - Inches(0.4),
             img_w, Inches(0.35), fill=RED)
    add_text(s6, ix, img_y + img_h - Inches(0.4),
             img_w, Inches(0.35),
             f'✗ {num}  {label}', size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_rect(s6, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.45),
         fill=NAVY_DARK)
add_text(s6, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.35),
         '反光 + 遮挡：目标特征被淹没',
         size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)

add_pagenum(s6, 6, total=8)


# ============================================================
# SLIDE 7: 失败案例 2 张大图（近命中 + 密集线）
# ============================================================
s7 = prs.slide_layouts[6]
s7 = prs.slides.add_slide(s7)
add_header(s7, '剩余难例 (2/3)',
           '近命中与密集线干扰案例')

img_y = Inches(1.45)
img_w = Inches(6.3)
img_h = Inches(5.4)
img_gap = Inches(0.3)
img_x_start = (SW - (img_w * 2 + img_gap)) // 2

fn_pairs = [
    ('23.png', '近命中（差 3px）', '#3'),
    ('d_57_.png', '密集线干扰', '#4'),
]

for i, (fn, label, num) in enumerate(fn_pairs):
    ix = img_x_start + i * (img_w + img_gap)
    add_rect(s7, ix, img_y, img_w, img_h, fill=BG_LIGHT)
    add_image_fit(s7, FN_DIR / fn,
                  ix + Inches(0.1), img_y + Inches(0.1),
                  img_w - Inches(0.2), img_h - Inches(0.5),
                  max_px=2400)
    add_rect(s7, ix, img_y + img_h - Inches(0.4),
             img_w, Inches(0.35), fill=RED)
    add_text(s7, ix, img_y + img_h - Inches(0.4),
             img_w, Inches(0.35),
             f'✗ {num}  {label}', size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_rect(s7, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.45),
         fill=NAVY_DARK)
add_text(s7, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.35),
         '预测接近但 IoU 刚好不达标 / 目标被前景线条遮挡',
         size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)

add_pagenum(s7, 7, total=8)


# ============================================================
# SLIDE 8: 失败案例 2 张大图（剩余 2 个）
# ============================================================
s8 = prs.slide_layouts[6]
s8 = prs.slides.add_slide(s8)
add_header(s8, '剩余难例 (3/3)',
           '反射叠加 + 特征极弱')

img_y = Inches(1.45)
img_w = Inches(6.3)
img_h = Inches(5.4)
img_gap = Inches(0.3)
img_x_start = (SW - (img_w * 2 + img_gap)) // 2

fn_pairs = [
    ('9.png', '反射 + 遮挡', '#5'),
    ('96.png', '特征极弱', '#6'),
]

for i, (fn, label, num) in enumerate(fn_pairs):
    ix = img_x_start + i * (img_w + img_gap)
    add_rect(s8, ix, img_y, img_w, img_h, fill=BG_LIGHT)
    add_image_fit(s8, FN_DIR / fn,
                  ix + Inches(0.1), img_y + Inches(0.1),
                  img_w - Inches(0.2), img_h - Inches(0.5),
                  max_px=2400)
    add_rect(s8, ix, img_y + img_h - Inches(0.4),
             img_w, Inches(0.35), fill=RED)
    add_text(s8, ix, img_y + img_h - Inches(0.4),
             img_w, Inches(0.35),
             f'✗ {num}  {label}', size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# 总结 banner
add_rect(s8, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.45),
         fill=NAVY_DARK)
add_text(s8, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.35),
         '6 个 FN 共占 13%，是数据本身的 hard limit',
         size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)

add_pagenum(s8, 8, total=8)


OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
print(f'✓ 已生成: {OUT}')
print(f'  大小: {OUT.stat().st_size / 1024:.1f} KB')
print(f'  页数: {len(prs.slides)}')